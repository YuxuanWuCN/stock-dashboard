# -*- coding: utf-8 -*-
"""scripts/rebuild_gold_deck.py —— 2026 中国国际大学生创新大赛 18 页金牌路演 PPT (全量高精构建脚本)

设计标准：
1. 16:9 宽屏 (13.333 x 7.5 英寸)
2. 科技金融专业色系：Navy (#0F172A)、Tech Blue (#2563EB)、Amber Gold (#D97706)、Emerald Green (#059669)、Crimson Red (#E11D48)
3. 结构化卡片容器、高光指标磁贴、加粗引导词、多列对比与学术图谱
4. 完美融合 TFAC 时变校准框架、2024-2026 年 300 标的大底座 ($t=3.92$)、绿电 (Sharpe 1.31 / MaxDD 12.8%) 与存储+黄金双资产杠铃回测
"""

from __future__ import annotations

import os
from pathlib import Path
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

C_NAVY = RGBColor(15, 23, 42)          # #0F172A 深海军蓝
C_BLUE = RGBColor(37, 99, 235)         # #2563EB 科技蓝
C_DARK_BLUE = RGBColor(30, 58, 138)    # #1E3A8A 暗蓝
C_GOLD = RGBColor(217, 119, 6)         # #D97706 琥珀金
C_GREEN = RGBColor(5, 150, 105)        # #059669 翡翠绿
C_RED = RGBColor(225, 29, 72)          # #E11D48 绯红
C_BG = RGBColor(248, 250, 252)         # #F8FAFC 浅灰白底色
C_CARD_BG = RGBColor(255, 255, 255)    # #FFFFFF 纯白卡片
C_CARD_BORDER = RGBColor(226, 232, 240)# #E2E8F0 浅灰边框
C_TEXT_MAIN = RGBColor(30, 41, 59)     # #1E293B 正文主色
C_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B 次要辅助字

FONT_NAME = "Microsoft YaHei"


def create_full_deck():
    prs = pptx.Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text: str, category_text: str = "2026 中国国际大学生创新大赛 · 达观数据产业命题赛道", page_num: int = 1):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.0), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_NAME
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_BLUE

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(10.5), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = FONT_NAME
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.color.rgb = C_NAVY

        pg_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.5), Inches(1.0), Inches(0.4))
        tf_pg = pg_box.text_frame
        tf_pg.margin_left = tf_pg.margin_top = tf_pg.margin_right = tf_pg.margin_bottom = 0
        p_pg = tf_pg.paragraphs[0]
        p_pg.alignment = PP_ALIGN.RIGHT
        p_pg.text = f"{page_num:02d} / 18"
        p_pg.font.name = FONT_NAME
        p_pg.font.size = Pt(12)
        p_pg.font.bold = True
        p_pg.font.color.rgb = C_TEXT_MUTED

    def add_card(slide, left: float, top: float, width: float, height: float, bg_color: RGBColor = C_CARD_BG, border_color: RGBColor = C_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.2)
        return shape

    # -------------------------------------------------------------------------
    # SLIDE 1: 封面页
    # -------------------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    add_card(s1, 0.8, 0.8, 11.733, 5.9, bg_color=RGBColor(241, 245, 249), border_color=C_BLUE)

    tbox = s1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(10.9), Inches(2.2))
    tf1 = tbox.text_frame
    tf1.word_wrap = True
    p0 = tf1.paragraphs[0]
    p0.text = "2026 中国国际大学生创新大赛 · 达观数据产业命题赛道"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_BLUE

    p1 = tf1.add_paragraph()
    p1.text = "Rainbow-FinGPT"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = C_NAVY

    p2 = tf1.add_paragraph()
    p2.text = "面向金融量化投研全流程的自主智能体系统"
    p2.font.name = FONT_NAME
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = C_DARK_BLUE

    p3 = tf1.add_paragraph()
    p3.text = "基于「定性语义 (FinEvidence) — 资产定价 (TFAC / Fama-MacBeth 3.0) — 战术风控 (Trend Gate)」三层解耦架构"
    p3.font.name = FONT_NAME
    p3.font.size = Pt(12)
    p3.font.color.rgb = C_TEXT_MUTED

    b_widths = 2.6
    b_lefts = [1.2, 4.0, 6.8, 9.6]
    b_titles = ["⚡ 自动化全闭环", "🧩 三层解耦架构", "🔍 100% 坐标级溯源", "📊 300 标的因果大底座"]
    b_subs = [
        "每日 18:00 自动抓取/定价/调仓\n耗时由数小时压至 15 分钟",
        "大模型做研报事实抽取\n定价与风控由纯数学公式驱动",
        "Citation-Grounded 段落锚定\n每条事实直通原文，消除幻觉",
        "2024~2026 跨周期 694 交易日\nHarvey Alpha t=3.92 ≥ 3.0"
    ]
    for i in range(4):
        add_card(s1, b_lefts[i], 3.7, b_widths, 1.8, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
        c_box = s1.shapes.add_textbox(Inches(b_lefts[i] + 0.15), Inches(3.8), Inches(b_widths - 0.3), Inches(1.6))
        c_tf = c_box.text_frame
        c_tf.word_wrap = True
        cp0 = c_tf.paragraphs[0]
        cp0.text = b_titles[i]
        cp0.font.name = FONT_NAME
        cp0.font.size = Pt(12)
        cp0.font.bold = True
        cp0.font.color.rgb = C_BLUE

        cp1 = c_tf.add_paragraph()
        cp1.text = b_subs[i]
        cp1.font.name = FONT_NAME
        cp1.font.size = Pt(9.5)
        cp1.font.color.rgb = C_TEXT_MAIN

    foot_box = s1.shapes.add_textbox(Inches(1.2), Inches(5.9), Inches(10.9), Inches(0.6))
    f_tf = foot_box.text_frame
    fp = f_tf.paragraphs[0]
    fp.text = "依托单位：华南师范大学阿伯丁数据科学与人工智能学院  |  团队负责人：吴宇轩  |  命题企业：达观数据有限公司"
    fp.font.name = FONT_NAME
    fp.font.size = Pt(11)
    fp.font.bold = True
    fp.font.color.rgb = C_NAVY

    # -------------------------------------------------------------------------
    # SLIDE 2: 24/7 研究助理
    # -------------------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "一个“永远不下班的 AI 投研助理” —— 破解传统投研效率与纪律困境", page_num=2)

    add_card(s2, 0.8, 1.5, 5.7, 5.3, bg_color=RGBColor(254, 242, 242), border_color=RGBColor(254, 202, 202))
    l_box = s2.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    ltf = l_box.text_frame
    ltf.word_wrap = True
    lp0 = ltf.paragraphs[0]
    lp0.text = "🔴 传统人工投研模式（痛点与瓶颈）"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(15)
    lp0.font.bold = True
    lp0.font.color.rgb = C_RED

    l_items = [
        ("繁琐低效：", "初级研究员 70% 精力被困在读研报、搬财报、抄表格，单篇研报复现需 4 ~ 20 小时；"),
        ("主观情绪：", "人工调仓极易受市场恐慌与贪婪干扰，追涨杀跌，缺乏严格因果纪律；"),
        ("经验断层：", "资深分析师定性认知无法代码化、可编程沉淀，行业覆盖度严重受限；"),
        ("时效滞后：", "宏观突发地缘或突发现货价格跳涨，人工往往在 3~5 个交易日后才能反应。")
    ]
    for bold_txt, norm_txt in l_items:
        p = ltf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + bold_txt
        r1.font.bold = True
        r1.font.color.rgb = C_RED
        r1.font.size = Pt(11)
        r2 = p.add_run()
        r2.text = norm_txt
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(11)

    add_card(s2, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    r_box = s2.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    rtf = r_box.text_frame
    rtf.word_wrap = True
    rp0 = rtf.paragraphs[0]
    rp0.text = "🟢 Rainbow-FinGPT 智能体全自动流水线"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(15)
    rp0.font.bold = True
    rp0.font.color.rgb = C_GREEN

    r_items = [
        ("24/7 自动闭环：", "交易日 18:00 自动触发，抓取、清洗、定价、风控、推送全部在 15 分钟内完成；"),
        ("节约 92% 劳动：", "将分析师从重复性数据搬运中彻底解放，端到端自动化交付出版级研报；"),
        ("三层解耦赋能：", "大模型负责研报因果事实抽取，定价与风控交由纯确定性数学公式；"),
        ("全额扣除摩擦：", "全流程扣除买 0.125% + 卖 0.175% 真实印花税与佣金，8% 死区控制换手。")
    ]
    for bold_txt, norm_txt in r_items:
        p = rtf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + bold_txt
        r1.font.bold = True
        r1.font.color.rgb = C_GREEN
        r1.font.size = Pt(11)
        r2 = p.add_run()
        r2.text = norm_txt
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(11)

    # -------------------------------------------------------------------------
    # SLIDE 3: 投研范式跃迁 (三方矩阵)
    # -------------------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "投研范式跃迁 —— 传统人工 vs 通用大模型炒股 vs Rainbow-FinGPT", page_num=3)

    add_card(s3, 0.8, 1.5, 11.733, 5.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    t_shape = s3.shapes.add_table(7, 4, Inches(1.0), Inches(1.7), Inches(11.333), Inches(4.8))
    table = t_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(2.8)
    table.columns[2].width = Inches(3.0)
    table.columns[3].width = Inches(3.333)

    headers = ["评价维度", "传统人工投研", "通用大模型直接预测 (如GPT/DeepSeek)", "Rainbow-FinGPT (本项目)"]
    for c_idx, h in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY if c_idx < 3 else C_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    matrix_rows = [
        ("核心运作模式", "人工读研报、搬数据、算表格", "直接让大模型读新闻给买卖买卖点", "大模型抽取客观事实 + 纯数学定价风控"),
        ("单篇研报耗时", "4 ~ 20 小时 / 篇", "秒级 (但结果完全不可信)", "约 15 分钟 (全流程工业自动闭环)"),
        ("财报数据可靠性", "人工复核，易漏看错看", "❌ 严重数值幻觉 (凭空捏造财报)", "🌟 100% 坐标级锚定，彻底消除幻觉"),
        ("决策因果可解释性", "分析师主观经验归因", "❌ 黑盒不可解释 (机构不敢跟投)", "🌟 资产定价与特质 Alpha 公式透明可溯"),
        ("时序前视未来函数", "受限于人力覆盖度", "❌ 严重时序穿越 (实盘亏损惨重)", "🌟 严格物理因果隔离 (仅用 <=t 日切片)"),
        ("战术风控与回撤", "依赖人工止损，易受情绪干扰", "❌ 缺乏战术门禁 (微观破位暴跌)", "🌟 Trend Gate™ C 浪硬门禁清仓防守")
    ]
    for r_idx, row in enumerate(matrix_rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(248, 250, 252) if c_idx < 3 else RGBColor(238, 242, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_NAME
            p.font.size = Pt(10)
            if c_idx == 3:
                p.font.bold = True
                p.font.color.rgb = C_BLUE
            else:
                p.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 4: 大模型三大死穴
    # -------------------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "为什么不能把通用大模型直接丢进二级市场？ —— 剖析三大致命死穴", page_num=4)

    c4_widths = 3.6
    c4_lefts = [0.8, 4.85, 8.9]
    c4_titles = ["❌ 死穴一：数值幻觉 (Hallucination)", "❌ 死穴二：黑盒决策 (Black-Box)", "❌ 死穴三：前视偏差 (Look-Ahead)"]
    c4_bolds = ["凭空捏造财报数据与利润指标", "端到端神经网络无法归因穿透", "预训练语料混杂未来时间信息"]
    c4_details = [
        "大模型本质是自回归概率生成，在存货减值、预付款比例等精密金融数值上频繁出错，无法通过投委会严谨审计。",
        "机构投资者掌管数十亿资金，不可能把资金托付给无法解释敞口来源、无法量化风险暴露的黑盒模型。",
        "大模型在训练时已经‘看’过了 2024~2026 年的行情结果，在回测中虚高数十倍，实盘上线瞬间一触即溃。"
    ]
    for i in range(3):
        add_card(s4, c4_lefts[i], 1.5, c4_widths, 4.3, bg_color=RGBColor(254, 242, 242), border_color=RGBColor(254, 202, 202))
        box = s4.shapes.add_textbox(Inches(c4_lefts[i] + 0.2), Inches(1.7), Inches(c4_widths - 0.4), Inches(3.9))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = c4_titles[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = C_RED

        p1 = tf.add_paragraph()
        p1.text = c4_bolds[i]
        p1.font.name = FONT_NAME
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY

        p2 = tf.add_paragraph()
        p2.text = c4_details[i]
        p2.font.name = FONT_NAME
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = C_TEXT_MAIN

    # 底部破局横条
    add_card(s4, 0.8, 6.0, 11.733, 0.9, bg_color=RGBColor(238, 242, 255), border_color=C_BLUE)
    sol_box = s4.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.7))
    stf = sol_box.text_frame
    stf.word_wrap = True
    sp = stf.paragraphs[0]
    sp.text = "💡 破局之道：三层解耦与物理时序隔离 —— 让大语言模型做语言理解（抽取事实），让数学公式与状态机掌控资金安全！"
    sp.font.name = FONT_NAME
    sp.font.size = Pt(11.5)
    sp.font.bold = True
    sp.font.color.rgb = C_BLUE

    # -------------------------------------------------------------------------
    # SLIDE 5: 拒绝因子动物园与学术诚信
    # -------------------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "坚守金融学本质 —— 拒绝“因子动物园”与暴力过拟合", page_num=5)

    add_card(s5, 0.8, 1.5, 5.7, 5.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    l5_box = s5.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l5_tf = l5_box.text_frame
    l5_tf.word_wrap = True
    lp0 = l5_tf.paragraphs[0]
    lp0.text = "⚠️ 传统量化陷阱：Factor Zoo 伪因子泛滥"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(14)
    lp0.font.bold = True
    lp0.font.color.rgb = C_GOLD

    items_l5 = [
        ("暴力挖掘：", "通过遗传规划挖掘数千个无经济学机理的高频公式，硬凑历史拟合度；"),
        ("多重检验失效：", "Harvey (2016) 顶刊指出传统 t>2.0 已严重失效，伪因子泛滥成灾；"),
        ("实盘即暴雷：", "过拟合策略上线即失效，遇到风格突变立即遭受巨大回撤。")
    ]
    for b, n in items_l5:
        p = l5_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_GOLD
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    add_card(s5, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    r5_box = s5.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r5_tf = r5_box.text_frame
    r5_tf.word_wrap = True
    rp0 = r5_tf.paragraphs[0]
    rp0.text = "🌟 Rainbow-FinGPT 的学术坚持与自律"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_GREEN

    items_r5 = [
        ("先验经济学驱动：", "仅构建具备坚实产业因果的因子（供需周期、克金成本、消纳率）；"),
        ("Fama-MacBeth 3.0：", "滚动两阶段截面回归，Newey-West HAC (q=4) 稳健估计；"),
        ("严苛 t ≥ 3.0 门禁：", "全市场 300 标的大底座实测 Harvey Alpha t = 3.92 (p < 0.01)；"),
        ("真实反例验证：", "立新能源暴涨 +82.36%，因特质 p=0.3543 被系统果断判定 REJECT 拦截！")
    ]
    for b, n in items_r5:
        p = r5_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_GREEN
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 6: 三层解耦架构拓扑
    # -------------------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "三层解耦总体架构 —— 定性认知 · 资产定价 · 战术风控", page_num=6)

    # 3 个分层大横条
    l_tops = [1.5, 3.2, 4.9]
    l_names = ["Layer 1 · 定性认知与事实抽取层 (Qualitative Cognition)", "Layer 2 · 资产定价与自适应校准层 (Asset Pricing & TFAC)", "Layer 3 · 战术风控与门控执行层 (Tactical Risk Gate)"]
    l_colors = [C_BLUE, C_DARK_BLUE, C_GREEN]
    l_descs = [
        "【FinEvidence 事实图谱引擎】: FOI 三元分离 (Fact/Opinion/Inference) + 100% 坐标级段落锚定 + 供应链卡位打分 (CS ≥ 12 门禁) ── 彻底消除幻觉",
        "【Fama-MacBeth 3.0 + TFAC 在线学习校准】: Carhart 4 因子 Newey-West 回归 + Hedge 算法在线更新 (遗憾界 O(√T)) + 二项显著性检验 + 拒绝预测 ── 动态捕获 Alpha",
        "【Trend Gate™ 布尔硬门禁 + 因果 ZigZag】: 斐波那契 [0.500, 0.618] 黄金支撑带右侧确认 + C 浪破位一票否决清仓 ── 强力压缩回撤"
    ]
    for i in range(3):
        add_card(s6, 0.8, l_tops[i], 11.733, 1.5, bg_color=C_CARD_BG, border_color=l_colors[i])
        box = s6.shapes.add_textbox(Inches(1.0), Inches(l_tops[i] + 0.15), Inches(11.333), Inches(1.2))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = l_names[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = l_colors[i]

        p1 = tf.add_paragraph()
        p1.text = l_descs[i]
        p1.font.name = FONT_NAME
        p1.font.size = Pt(10.5)
        p1.font.color.rgb = C_TEXT_MAIN

    # 底部说明
    add_card(s6, 0.8, 6.55, 11.733, 0.5, bg_color=RGBColor(241, 245, 249), border_color=C_CARD_BORDER)
    f6_box = s6.shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.333), Inches(0.4))
    f6_tf = f6_box.text_frame
    p = f6_tf.paragraphs[0]
    p.text = "🔒 核心公理：大模型定性归纳与数值计算绝对物理隔离，大模型绝不碰浮点技术指标或仓位分配！"
    p.font.name = FONT_NAME
    p.font.size = Pt(10.5)
    p.font.bold = True
    p.font.color.rgb = C_NAVY

    # -------------------------------------------------------------------------
    # SLIDE 7: Layer 1 FinEvidence
    # -------------------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Layer 1 · FinEvidence 研报因果事实图谱抽取器 —— 100% 坐标级段落溯源", page_num=7)

    c7_w = 3.6
    c7_lefts = [0.8, 4.85, 8.9]
    c7_titles = ["1. FOI 三元分离机制", "2. 100% 坐标级段落锚定", "3. 供应链卡位评分 (CS)"]
    c7_points = [
        [("严格隔离：", "将研报拆解为 [FACT] 财务真实值、[OPINION] 卖方主观预期与 [INFERENCE] 逻辑推论；"), ("防偏见污染：", "客观财务数据进入定价矩阵，主观过度乐观预期进入风险折价。")],
        [("Citation-Grounded：", "每条提取结论强制绑定 PDF 原件页码与段落坐标（如《XX研报》P.14 第3段）；"), ("可穿透审计：", "支持监管与投资委员会一键穿透核查，彻底攻克大模型数值幻觉。")],
        [("10 题卡位矩阵：", "考察先进制程占比、自研主控率、AISC 克金成本、特高压消纳率；"), ("硬门槛准入：", "设定 CS ≥ 12 分硬门禁，过滤蹭概念杂毛股，锁定真正核心龙头。")]
    ]
    for i in range(3):
        add_card(s7, c7_lefts[i], 1.5, c7_w, 5.3, bg_color=C_CARD_BG, border_color=C_BLUE)
        box = s7.shapes.add_textbox(Inches(c7_lefts[i] + 0.2), Inches(1.7), Inches(c7_w - 0.4), Inches(4.9))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = c7_titles[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = C_BLUE

        for b, n in c7_points[i]:
            p = tf.add_paragraph()
            r1 = p.add_run()
            r1.text = "• " + b
            r1.font.bold = True
            r1.font.color.rgb = C_NAVY
            r2 = p.add_run()
            r2.text = n
            r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 8: Layer 2 Fama-MacBeth 3.0 & TFAC
    # -------------------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Layer 2 · Fama-MacBeth 3.0 与 TFAC 在线学习时变校准框架", page_num=8)

    add_card(s8, 0.8, 1.5, 5.7, 5.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    l8_box = s8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l8_tf = l8_box.text_frame
    l8_tf.word_wrap = True
    lp0 = l8_tf.paragraphs[0]
    lp0.text = "📐 Fama-MacBeth 3.0 截面资产定价"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(14)
    lp0.font.bold = True
    lp0.font.color.rgb = C_NAVY

    items_l8 = [
        ("两阶段滚动回归：", "基于 Carhart 4 因子（MKT, SMB, HML, MOM）在 252 日滚动窗口中剥离系统性风格暴露；"),
        ("Newey-West HAC：", "自适应最优滞后阶数 q = 4，解决金融时序自相关与异方差问题；"),
        ("特质 Alpha 提取：", "仅保留 p < 0.05 且信息比率 IR ≥ 0.30 的真实超额标的。")
    ]
    for b, n in items_l8:
        p = l8_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    add_card(s8, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(238, 242, 255), border_color=C_BLUE)
    r8_box = s8.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r8_tf = r8_box.text_frame
    r8_tf.word_wrap = True
    rp0 = r8_tf.paragraphs[0]
    rp0.text = "🚀 TFAC 在线学习时变校准框架 (学术创新)"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_BLUE

    items_r8 = [
        ("Hedge 算法在线更新：", "将在线学习引入因子方向校准，理论证明累积遗憾界 O(√T ln K)；"),
        ("单侧二项显著性检验：", "以置信度 ≥ 70% 作为因子有效门槛，杜绝过拟合经验阈值；"),
        ("主动拒绝预测 (Reject Option)：", "震荡无把握期输出 INVALID 主动持币观望；"),
        ("实证显著提升：", "有效方向胜率提升至 57.60%，策略夏普比率提升至 1.31！")
    ]
    for b, n in items_r8:
        p = r8_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_BLUE
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 9: Layer 2 NALE 产业链图谱
    # -------------------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Layer 2 · NALE 产业链拓扑图与高频现货传导 —— 领先卖方 5 日捕获脉冲", page_num=9)

    add_card(s9, 0.8, 1.5, 11.733, 2.4, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    t9_box = s9.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.333), Inches(2.0))
    t9_tf = t9_box.text_frame
    t9_tf.word_wrap = True
    p0 = t9_tf.paragraphs[0]
    p0.text = "🌐 五级产业链拓扑图谱与高频现货直连网络"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = C_NAVY

    p1 = t9_tf.add_paragraph()
    p1.text = "【上游材料/设备】 ──> 【晶圆制造/代工】 ──> 【芯片设计/主控】 ──> 【模组封装测试】 ──> 【终端集成应用】\n• 存储高频直连：接入 TrendForce 现货 DXI 指数、韩国海关芯片出口月报、预付款与库存周转数据\n• 黄金高频直连：接入上海金交所 Au99.99、COMEX 黄金期货、美国实际利率与全球央行购金储备"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(11)
    p1.font.color.rgb = C_TEXT_MAIN

    add_card(s9, 0.8, 4.1, 11.733, 2.7, bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    b9_box = s9.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.333), Inches(2.3))
    b9_tf = b9_box.text_frame
    b9_tf.word_wrap = True
    bp0 = b9_tf.paragraphs[0]
    bp0.text = "⚡ 网络阻尼传导方程与时效优势实证"
    bp0.font.name = FONT_NAME
    bp0.font.size = Pt(14)
    bp0.font.bold = True
    bp0.font.color.rgb = C_GREEN

    bp1 = b9_tf.add_paragraph()
    bp1.text = "• 阻尼传导公式：S_NALE = (1 - α) * S_0 + α * (W * S_0)，阻尼系数 α = 0.4，精准模拟价格脉冲沿产业链上下游的衰减与扩散；\n• 领先卖方 5 日：卖方研报从调研到发布通常需要 3~7 天，NALE 引擎在高频现货异动当日即完成全产业链得分重构；\n• 锁价优势捕捉：在 2025 年存储超级周期中，提前捕捉到佰维存储、德明利的低价晶圆库存与主控毛利修复，吃满主升浪。"
    bp1.font.name = FONT_NAME
    bp1.font.size = Pt(11)
    bp1.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 10: Layer 3 量化买卖决策机制与状态机触发方程 (什么时候买？什么时候卖？)
    # -------------------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Layer 3 · 量化买卖决策机制与状态机触发方程 —— 什么时候买？什么时候卖？", page_num=10)

    # 左侧卡片：买入判定
    add_card(s10, 0.8, 1.5, 5.7, 5.3, bg_color=RGBColor(240, 253, 244), border_color=C_GREEN)
    l10_box = s10.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l10_tf = l10_box.text_frame
    l10_tf.word_wrap = True
    lp0 = l10_tf.paragraphs[0]
    lp0.text = "🟢 买入 / 加仓触发机制（三重因果信号共振）"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(13.5)
    lp0.font.bold = True
    lp0.font.color.rgb = C_GREEN

    items_l10 = [
        ("1. 基本面与研报门禁 (Layer 1)：", "FinEvidence 抽取供应链卡位分 CS ≥ 12 且无重大财务/减值风险预警；"),
        ("2. TFAC 定价与显著性 (Layer 2)：", "TFAC 在线校准输出 POSITIVE/LONG，二项显著性置信度 ≥ 70%（非 INVALID），特质 Alpha p < 0.05 且 IR ≥ 0.30；"),
        ("3. 战术择时与量价共振 (Layer 3)：", "股价站上均线 (Price > MA20)，MACD > 0，且因果 ZigZag 处于突破期或回踩斐波那契 [0.500, 0.618] 支撑带且缩量 ≥ 20%。"),
        ("⚡ 买入判定方程：", "EntryPass = (CS ≥ 12) & (TFAC_Dir == LONG) & (Conf ≥ 70%) & (Price > MA20) & !(Phase == Phase_C)")
    ]
    for b, n in items_l10:
        p = l10_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY if "方程" not in b else C_GREEN
        r1.font.size = Pt(10)
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(9.5)

    # 右侧卡片：卖出判定
    add_card(s10, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(254, 242, 242), border_color=C_RED)
    r10_box = s10.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r10_tf = r10_box.text_frame
    r10_tf.word_wrap = True
    rp0 = r10_tf.paragraphs[0]
    rp0.text = "🔴 卖出 / 减仓 / 清仓机制（四重硬门禁与止损）"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(13.5)
    rp0.font.bold = True
    rp0.font.color.rgb = C_RED

    items_r10 = [
        ("1. Trend Gate 破位清仓：", "股价跌破 MA20 且状态机确认进入 C 浪主跌破位期，系统强制 100% 清仓置为 CASH，拒绝死扛；"),
        ("2. TFAC 因子失效与反转：", "TFAC 二项置信度跌破 70%（输出 INVALID 拒绝预测）或方向反转为 SHORT，主动降仓或清仓观望；"),
        ("3. 基本面减值冲击：", "上游现货 DXI 价格暴跌或海关出口同比负增长触发减值惩罚，GFCA 分数跌出 Top 3 换仓退出；"),
        ("4. 动态追踪止损：", "单笔浮亏达到动态 ATR 止损线（-5%~-8%）或高位顶背离强制止盈。"),
        ("🛡️ 清仓判定方程：", "ExitPass = (Price < MA20 & Phase == Phase_C) | (TFAC_Conf < 70%) | (Loss > ATR_Stop)")
    ]
    for b, n in items_r10:
        p = r10_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY if "方程" not in b else C_RED
        r1.font.size = Pt(10)
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(9.5)

    # -------------------------------------------------------------------------
    # SLIDE 11: 每日 18:00 自动闭环
    # -------------------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "每日 18:00 投研全流程无人值守自动闭环 —— 7 大工业级流水线", page_num=11)

    steps_w = 1.5
    steps_lefts = [0.8 + i * 1.7 for i in range(7)]
    step_nums = ["01", "02", "03", "04", "05", "06", "07"]
    step_names = ["数据抓取", "多源清洗", "事实抽取", "因子定价", "战术门控", "调仓结算", "报告推送"]
    step_descs = [
        "18:00 收盘后自动拉取行情/现货",
        "清洗缺失值，对齐 CSMAR 因子库",
        "FinEvidence 解析最新研报与财报",
        "TFAC 在线校准与 Fama-MacBeth 回归",
        "Trend Gate 评估 C 浪与 MA20 状态",
        "8% 死区调仓，全额计提佣金印花税",
        "自动生成 300 DPI 研报并邮件推送"
    ]
    for i in range(7):
        add_card(s11, steps_lefts[i], 1.5, steps_w, 4.3, bg_color=C_CARD_BG, border_color=C_BLUE)
        box = s11.shapes.add_textbox(Inches(steps_lefts[i] + 0.1), Inches(1.7), Inches(steps_w - 0.2), Inches(3.9))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.text = step_nums[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = C_BLUE

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.text = step_names[i]
        p1.font.name = FONT_NAME
        p1.font.size = Pt(12)
        p1.font.bold = True
        p1.font.color.rgb = C_NAVY

        p2 = tf.add_paragraph()
        p2.text = step_descs[i]
        p2.font.name = FONT_NAME
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = C_TEXT_MAIN

    add_card(s11, 0.8, 6.0, 11.733, 0.9, bg_color=RGBColor(240, 253, 244), border_color=RGBColor(187, 247, 208))
    t11_box = s11.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.7))
    tf11 = t11_box.text_frame
    p = tf11.paragraphs[0]
    p.text = "⏱️ 端到端耗时：由人工的 4~20 小时压缩至 15 分钟以内，节约 92% 重复性劳动，真正做到“闭市即交付”！"
    p.font.name = FONT_NAME
    p.font.size = Pt(11.5)
    p.font.bold = True
    p.font.color.rgb = C_GREEN

    # -------------------------------------------------------------------------
    # SLIDE 12: 工业级工程底座与质量门禁
    # -------------------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "工业级工程底座与质量门禁 —— 可审计、防篡改、无缝迁移", page_num=12)

    c12_w = 5.7
    add_card(s12, 0.8, 1.5, c12_w, 5.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    l12_box = s12.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l12_tf = l12_box.text_frame
    l12_tf.word_wrap = True
    lp0 = l12_tf.paragraphs[0]
    lp0.text = "🛡️ 自动化质量门禁与测试体系"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(14)
    lp0.font.bold = True
    lp0.font.color.rgb = C_NAVY

    items_l12 = [
        ("78 项单元测试 100% 通过：", "覆盖状态机、TFAC 因子校准、调仓死区与因子契约，持续集成保护；"),
        ("SHA-256 数据防篡改指纹：", "每份研报与回测结果均绑定数据哈希指纹，防止人为后验修改；"),
        ("变异测试与弱断言拦截：", "拒绝假测试，确保每个关键数学逻辑均有可手算独立基准复核。")
    ]
    for b, n in items_l12:
        p = l12_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    add_card(s12, 6.8, 1.5, c12_w, 5.3, bg_color=C_CARD_BG, border_color=C_BLUE)
    r12_box = s12.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r12_tf = r12_box.text_frame
    r12_tf.word_wrap = True
    rp0 = r12_tf.paragraphs[0]
    rp0.text = "🔌 高校 CSMAR 因子库契约与无缝迁移"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_BLUE

    items_r12 = [
        ("官方 CSMAR 因子对齐：", "严格遵循国泰安 6 列标准规范（date, MKT, SMB, HML, MOM, rf），数据契约 100% 校验；"),
        ("出版级 PDF 渲染引擎：", "基于 ReportLab 自动编译 300 DPI 矢量图表，对齐国际顶刊排版规范；"),
        ("零成本迁移至券商中台：", "标准化 Python 接口与微服务容器化封装，支持私有化部署。")
    ]
    for b, n in items_r12:
        p = r12_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_BLUE
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 13: 实证一 · 存储超级周期
    # -------------------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "实证一 · A股半导体存储超级周期实测 —— 吃满主升浪与高位逃顶", page_num=13)

    m_w = 3.6
    m_lefts = [0.8, 4.85, 8.9]
    m_vals = ["+267.35%", "2.51 / 4.63", "29.14% / 6.51%"]
    m_titles = ["策略累计收益 (年化 +218%)", "年化夏普比率 (极高胜率)", "最大回撤 (避开腰斩暴跌)"]
    m_colors = [C_GOLD, C_BLUE, C_GREEN]
    for i in range(3):
        add_card(s13, m_lefts[i], 1.5, m_w, 1.8, bg_color=C_CARD_BG, border_color=m_colors[i])
        box = s13.shapes.add_textbox(Inches(m_lefts[i] + 0.1), Inches(1.6), Inches(m_w - 0.2), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.text = m_vals[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(26)
        p0.font.bold = True
        p0.font.color.rgb = m_colors[i]

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.text = m_titles[i]
        p1.font.name = FONT_NAME
        p1.font.size = Pt(10.5)
        p1.font.color.rgb = C_TEXT_MAIN

    add_card(s13, 0.8, 3.5, 11.733, 3.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    t13_shape = s13.shapes.add_table(5, 5, Inches(1.0), Inches(3.7), Inches(11.333), Inches(2.9))
    t13 = t13_shape.table
    for c_idx, w in enumerate([Inches(2.5), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.233)]):
        t13.columns[c_idx].width = w

    t13_headers = ["策略方案 / 标的", "累计收益率", "年化收益率", "夏普比率 (Sharpe)", "最大回撤 (MaxDD)"]
    for c_idx, h in enumerate(t13_headers):
        cell = t13.cell(0, c_idx)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_NAVY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_NAME
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    rows13 = [
        ("👑 Rainbow-FinGPT 存储策略", "+267.35%", "+218.23%", "2.51 (卡尔玛 7.49)", "29.14% (锁住利润)"),
        ("⚡ 存储超级周期专用精测", "+159.01%", "+105.65%", "4.63 (卡尔玛 16.22)", "6.51% (极低回撤)"),
        ("📊 存储 5 巨头等权基准", "+159.20%", "+128.50%", "1.48", "-54.13% (深度腰斩)"),
        ("📉 半导体芯片 ETF (512760.SH)", "+98.90%", "+76.40%", "1.12", "-38.20%")
    ]
    for r_idx, row in enumerate(rows13):
        for c_idx, val in enumerate(row):
            cell = t13.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(238, 242, 255) if r_idx < 2 else RGBColor(248, 250, 252)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_NAME
            p.font.size = Pt(10)
            if r_idx < 2:
                p.font.bold = True
                p.font.color.rgb = C_BLUE
            else:
                p.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 14: 实证二与三 · 黄金避险与绿电实测
    # -------------------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "实证二与三 · 黄金地缘避险与绿电公用事业实测 —— 跨行业稳健对冲", page_num=14)

    add_card(s14, 0.8, 1.5, 5.7, 5.3, bg_color=C_CARD_BG, border_color=C_GOLD)
    l14_box = s14.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l14_tf = l14_box.text_frame
    l14_tf.word_wrap = True
    lp0 = l14_tf.paragraphs[0]
    lp0.text = "🪙 实证二：黄金贵金属地缘避险专题"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(14)
    lp0.font.bold = True
    lp0.font.color.rgb = C_GOLD

    items_l14 = [
        ("策略累计收益：", "+94.84% (年化 +105.82%)，同期黄金 ETF 仅 +28.22%；"),
        ("年化夏普比率：", "1.67 (卡尔玛比率 5.62)，超额 Alpha 斩获 +66.62%；"),
        ("最大动态回撤：", "18.82% (相比单股最大跌幅压降超 15 个百分点)；"),
        ("核心选股 Alpha：", "重仓紫金矿业、山东黄金、赤峰黄金，精准捕捉央行购金与地缘脉冲。")
    ]
    for b, n in items_l14:
        p = l14_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_GOLD
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    add_card(s14, 6.8, 1.5, 5.7, 5.3, bg_color=C_CARD_BG, border_color=C_GREEN)
    r14_box = s14.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r14_tf = r14_box.text_frame
    r14_tf.word_wrap = True
    rp0 = r14_tf.paragraphs[0]
    rp0.text = "⚡ 实证三：绿电公用事业板块终审实测"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_GREEN

    items_r14 = [
        ("策略夏普比率：", "1.31 (经状态机与动态仓位优化提升 10.1%)；"),
        ("最大动态回撤：", "12.80% (相比基准 -33.05% 强力压降 61.3%)；"),
        ("策略累计收益：", "+56.09% (年化 +30.20%)，卡尔玛比率 2.36；"),
        ("消纳率因子赋能：", "筛选长江电力、三峡能源等具备特高压消纳优势的龙头，平稳防御。")
    ]
    for b, n in items_r14:
        p = r14_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_GREEN
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 15: 严谨样本外检验 (OOS Walk-Forward) 与双资产杠铃敏感性实证 (拒绝过拟合)
    # -------------------------------------------------------------------------
    s15 = prs.slides.add_slide(blank_layout)
    add_header(s15, "实证四 · 样本外向前行走检验 (OOS Walk-Forward) 与双资产杠铃分析", page_num=15)

    # 左侧：Walk-Forward 与 300 标的大底座
    add_card(s15, 0.8, 1.5, 5.7, 5.3, bg_color=C_CARD_BG, border_color=C_BLUE)
    l15_box = s15.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l15_tf = l15_box.text_frame
    l15_tf.word_wrap = True
    lp0 = l15_tf.paragraphs[0]
    lp0.text = "🌐 样本外向前行走检验 (OOS Walk-Forward)"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(13)
    lp0.font.bold = True
    lp0.font.color.rgb = C_BLUE

    items_l15 = [
        ("严格时序隔离：", "样本内 (IS: 180日) 定参；样本外 (OOS: 514日) 完全前向盲测，物理隔离无穿越；"),
        ("300 标的 69,300 独立样本：", "覆盖全市场 300 标的，样本外 Harvey Alpha t = 3.92 ≥ 3.0 (p < 0.01)；"),
        ("方向命中与校准：", "5日样本外方向命中率 53.32%，Brier 概率校准得分 0.2665；"),
        ("宽平坦参数盆地 (Flat Basin)：", "滚动回看窗口 H ∈ [15, 60] 天内策略夏普稳定在 1.18~1.31（非过拟合孤岛）。")
    ]
    for b, n in items_l15:
        p = l15_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY
        r1.font.size = Pt(9.8)
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(9.3)

    # 右侧：存储+黄金双资产动态杠铃与极端压力测试
    add_card(s15, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(240, 253, 244), border_color=C_GREEN)
    r15_box = s15.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r15_tf = r15_box.text_frame
    r15_tf.word_wrap = True
    rp0 = r15_tf.paragraphs[0]
    rp0.text = "⚖️ 存储 + 黄金双资产杠铃与黑天鹅压力测试"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(13)
    rp0.font.bold = True
    rp0.font.color.rgb = C_GREEN

    items_r15 = [
        ("近乎正交弱相关性：", "存储 vs 黄金日收益率相关系数 ρ = 0.0768，天然跨周期避险对冲；"),
        ("多资产分散化增益：", "分散化比率 DR = 1.42，组合波动率相比单一存储板块大幅降低 40.8%；"),
        ("动态状态机切换：", "牛市 80% 存储进攻，熊市 80% 黄金避险，最大回撤压制至 12.50%；"),
        ("极端黑天鹅压力测试：", "在 2024 年初雪球爆仓极端暴跌期，系统主动拒绝率达 78.5% (持币防守)，将最大回撤由 -18.4% 压制至 -4.2%！")
    ]
    for b, n in items_r15:
        p = r15_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_NAVY
        r1.font.size = Pt(9.8)
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN
        r2.font.size = Pt(9.3)

    # -------------------------------------------------------------------------
    # SLIDE 16: 达观数据 10 项考核满分总成绩单
    # -------------------------------------------------------------------------
    s16 = prs.slides.add_slide(blank_layout)
    add_header(s16, "达观数据 10 项核心考核指标达标总成绩单 —— 全线 100% 超额达成", page_num=16)

    kpi_w = 2.6
    kpi_lefts = [0.8, 3.84, 6.88, 9.93]
    kpi_vals = ["92.4%", "98.9%", "100.0%", "+218.2%"]
    kpi_labels = ["研报提取正确率 (门槛 ≥80%)", "代码成功运行率 (门槛 ≥90%)", "证据可追溯覆盖率 (门槛 ≥95%)", "最高年化超额收益 (门槛 ≥10%)"]
    kpi_colors = [C_BLUE, C_GREEN, C_DARK_BLUE, C_GOLD]
    for i in range(4):
        add_card(s16, kpi_lefts[i], 1.5, kpi_w, 1.8, bg_color=C_CARD_BG, border_color=kpi_colors[i])
        box = s16.shapes.add_textbox(Inches(kpi_lefts[i] + 0.1), Inches(1.6), Inches(kpi_w - 0.2), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        p0.text = kpi_vals[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(28)
        p0.font.bold = True
        p0.font.color.rgb = kpi_colors[i]

        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        p1.text = kpi_labels[i]
        p1.font.name = FONT_NAME
        p1.font.size = Pt(9.5)
        p1.font.color.rgb = C_TEXT_MAIN

    # 10 项考核清单卡片
    add_card(s16, 0.8, 3.5, 11.733, 3.3, bg_color=C_CARD_BG, border_color=C_CARD_BORDER)
    box16 = s16.shapes.add_textbox(Inches(1.0), Inches(3.7), Inches(11.333), Inches(2.9))
    tf16 = box16.text_frame
    tf16.word_wrap = True
    p0 = tf16.paragraphs[0]
    p0.text = "📋 产业命题 10 项核心要求逐项对照验收结论"
    p0.font.name = FONT_NAME
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = C_NAVY

    p1 = tf16.add_paragraph()
    p1.text = "1. 研报数据结构化抽取：FOI 三元分离，正确率 92.4% (超额达标)\n2. 财务关键指标对齐：Wind / CSMAR 官方标准 100% 对齐\n3. 多因子资产定价模型：Fama-MacBeth 3.0 + TFAC 在线校准，Harvey t=3.92 (超额达标)\n4. 供应链网络知识图谱：NALE 五级拓扑网络阻尼传导，领先卖方 5 日 (超额达标)\n5. 战术风控与回撤控制：Trend Gate C 浪硬门禁，绿电回撤降至 12.80% (超额达标)\n6. 自动化流程执行：每日 18:00 任务调度，15 分钟全自动交付 (超额达标)\n7. 证据链可解释性：Citation-Grounded 坐标级段落绑定 100% (超额达标)\n8. 策略投资表现：存储 +267.35%，黄金 +94.84%，绿电 Sharpe 1.31 (全线超额达成)"
    p1.font.name = FONT_NAME
    p1.font.size = Pt(10)
    p1.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 17: 学术诚信案例 (立新能源)
    # -------------------------------------------------------------------------
    s17 = prs.slides.add_slide(blank_layout)
    add_header(s17, "学术诚信与风控边界 —— 涨了 +82.36%，系统为何依然果断拒绝？", page_num=17)

    add_card(s17, 0.8, 1.5, 5.7, 5.3, bg_color=RGBColor(254, 242, 242), border_color=C_RED)
    l17_box = s17.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(5.3), Inches(4.9))
    l17_tf = l17_box.text_frame
    l17_tf.word_wrap = True
    lp0 = l17_tf.paragraphs[0]
    lp0.text = "🔍 现象：立新能源 (001258) 短期暴涨 +82.36%"
    lp0.font.name = FONT_NAME
    lp0.font.size = Pt(14)
    lp0.font.bold = True
    lp0.font.color.rgb = C_RED

    items_l17 = [
        ("市场盲目追高：", "样本期内立新能源因题材炒作累计大涨 +82.36%，散户与传统动量模型争相追涨；"),
        ("Fama-MacBeth 计量检验：", "系统剥离市场 Beta 与风格因子后，其特质 Alpha 检验 p = 0.3543 (远未达到 p < 0.05 门槛)；"),
        ("信息比率极低：", "特质信息比率 IR = 0.063，远低于 0.30 准入门槛。")
    ]
    for b, n in items_l17:
        p = l17_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_RED
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    add_card(s17, 6.8, 1.5, 5.7, 5.3, bg_color=RGBColor(240, 253, 244), border_color=C_GREEN)
    r17_box = s17.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.3), Inches(4.9))
    r17_tf = r17_box.text_frame
    r17_tf.word_wrap = True
    rp0 = r17_tf.paragraphs[0]
    rp0.text = "🛡️ 判定：系统果断执行 REJECT 拦截"
    rp0.font.name = FONT_NAME
    rp0.font.size = Pt(14)
    rp0.font.bold = True
    rp0.font.color.rgb = C_GREEN

    items_r17 = [
        ("一票否决拦截：", "系统判定其上涨为纯市场 Beta 漂移，缺乏核心特质因果，果断执行拒绝入池；"),
        ("成功避开退潮暴跌：", "该股随后在退潮期遭遇断崖式大跌，系统因提前拦截避免了大幅回撤；"),
        ("学术诚信底线：", "宁可错过短期泡沫，也绝不放宽计量门禁！用严密的学术纪律保护投资者本金安全。")
    ]
    for b, n in items_r17:
        p = r17_tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "• " + b
        r1.font.bold = True
        r1.font.color.rgb = C_GREEN
        r2 = p.add_run()
        r2.text = n
        r2.font.color.rgb = C_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 18: 终章总结与产教协同
    # -------------------------------------------------------------------------
    s18 = prs.slides.add_slide(blank_layout)
    add_header(s18, "终章总结 · 产教协同重塑量化投研新生态", page_num=18)

    c18_w = 3.6
    c18_lefts = [0.8, 4.85, 8.9]
    c18_titles = ["1. 产教融合标杆答卷", "2. 赋能达观大模型生态", "3. 华师阿伯丁团队力量"]
    c18_points = [
        [("达观命题 100% 达成：", "首创三层解耦与 TFAC 在线校准，实现理论与工业级工程双重突破；"), ("完整知识产权：", "沉淀白皮书、理论附录、代码库与实证研报全套学术成果。")],
        [("垂直插件赋能：", "作为量化投研中台插件，无缝嵌入达观‘曹植大模型’，赋能金融机构；"), ("降低 90% 成本：", "助力券商与基金公司将单篇研报复现成本由数小时压缩至 15 分钟。")],
        [("跨学科交叉：", "数据科学、金融工程与人工智能交叉融合，具备极强的代码自研与学术推导能力；"), ("持续迭代演进：", "未来将接入更多大宗商品与宏观因子，打造全天候智能投资中枢。")]
    ]
    for i in range(3):
        add_card(s18, c18_lefts[i], 1.5, c18_w, 4.3, bg_color=C_CARD_BG, border_color=C_BLUE)
        box = s18.shapes.add_textbox(Inches(c18_lefts[i] + 0.2), Inches(1.7), Inches(c18_w - 0.4), Inches(3.9))
        tf = box.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = c18_titles[i]
        p0.font.name = FONT_NAME
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = C_BLUE

        for b, n in c18_points[i]:
            p = tf.add_paragraph()
            r1 = p.add_run()
            r1.text = "• " + b
            r1.font.bold = True
            r1.font.color.rgb = C_NAVY
            r2 = p.add_run()
            r2.text = n
            r2.font.color.rgb = C_TEXT_MAIN

    # 底部致谢
    add_card(s18, 0.8, 6.0, 11.733, 0.9, bg_color=RGBColor(238, 242, 255), border_color=C_BLUE)
    end_box = s18.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.7))
    etf = end_box.text_frame
    ep = etf.paragraphs[0]
    ep.alignment = PP_ALIGN.CENTER
    ep.text = "✨ 感谢各位评委老师的倾听与指导！Rainbow-FinGPT 致力于让金融量化更严谨、更高效、更安全！"
    ep.font.name = FONT_NAME
    ep.font.size = Pt(13)
    ep.font.bold = True
    ep.font.color.rgb = C_BLUE

    # 保存
    out_path = Path("2026中国国际大学生创新大赛_网评路演PPT_Rainbow-FinGPT.pptx")
    prs.save(str(out_path))
    print(f"Successfully generated 18-slide gold deck to {out_path}")

    # 镜像同步
    mirror_path = Path("Rainbow_FinGPTv2/2026中国国际大学生创新大赛_网评路演PPT_Rainbow-FinGPT.pptx")
    mirror_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(mirror_path))
    print(f"Successfully mirrored deck to {mirror_path}")


if __name__ == "__main__":
    create_full_deck()
